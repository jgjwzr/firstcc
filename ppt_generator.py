"""
Stage 1: Extract images + translations from PDFs
Stage 2: Build PPT from reviewed translations
Usage:
  Stage 1: python ppt_generator.py stage1
  Stage 2: python ppt_generator.py stage2
"""

import fitz
import json
import os
import sys
import re
from pathlib import Path
from io import BytesIO
from deep_translator import GoogleTranslator
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# Configuration
# ============================================================
MAIN_PDF = "D:/firstcc/论文正文.pdf"
SI_PDF = "D:/firstcc/论文si.pdf"
IMAGES_DIR = "D:/firstcc/images"
TRANSLATION_FILE = "D:/firstcc/translations_draft.json"
OUTPUT_PPT = "D:/firstcc/output.pptx"
RENDER_DPI = 250

# ============================================================
# Scientific Glossary (EN -> CN)
# ============================================================
GLOSSARY = {
    "hydrogen evolution reaction": "析氢反应",
    "HER": "HER",
    "alkaline water electrolysis": "碱性水电解",
    "ALKWE": "ALKWE",
    "overpotential": "过电位",
    "current density": "电流密度",
    "Tafel slope": "Tafel斜率",
    "electrochemical impedance spectroscopy": "电化学阻抗谱",
    "EIS": "EIS",
    "electrochemically active surface area": "电化学活性面积",
    "ECSA": "ECSA",
    "double-layer capacitance": "双电层电容",
    "Faradaic efficiency": "法拉第效率",
    "linear sweep voltammetry": "线性扫描伏安法",
    "LSV": "LSV",
    "chronopotentiometry": "计时电位法",
    "scanning electron microscopy": "扫描电子显微镜",
    "SEM": "SEM",
    "transmission electron microscopy": "透射电子显微镜",
    "TEM": "TEM",
    "high-angle annular dark-field scanning TEM": "高角环形暗场扫描透射电子显微镜",
    "HAADF-STEM": "HAADF-STEM",
    "high-resolution TEM": "高分辨透射电子显微镜",
    "HRTEM": "HRTEM",
    "X-ray diffraction": "X射线衍射",
    "XRD": "XRD",
    "X-ray photoelectron spectroscopy": "X射线光电子能谱",
    "XPS": "XPS",
    "X-ray absorption near-edge structure": "X射线吸收近边结构",
    "XANES": "XANES",
    "extended X-ray absorption fine structure": "扩展X射线吸收精细结构",
    "EXAFS": "EXAFS",
    "Fourier-transformed EXAFS": "傅里叶变换EXAFS",
    "wavelet-transformed EXAFS": "小波变换EXAFS",
    "density functional theory": "密度泛函理论",
    "DFT": "DFT",
    "Gibbs free energy": "吉布斯自由能",
    "projected density of states": "投影态密度",
    "PDOS": "PDOS",
    "differential charge": "差分电荷",
    "mass transfer": "质量传输",
    "mass transport": "质量传输",
    "electron transfer": "电子转移",
    "intrinsic activity": "本征活性",
    "hierarchical porosity": "多级孔结构",
    "triscale porosity": "三尺度孔结构",
    "nano-micro-macro": "纳-微-宏",
    "hydrophilicity": "亲水性",
    "aerophobicity": "疏气性",
    "superaerophobicity": "超疏气性",
    "bubble detachment": "气泡脱附",
    "bubble evolution": "气泡演化",
    "adhesion force": "粘附力",
    "contact angle": "接触角",
    "electrolyte permeation": "电解质渗透",
    "powder metallurgy": "粉末冶金",
    "impregnation": "浸渍",
    "calcination": "煅烧",
    "polyoxometalate": "多金属氧酸盐",
    "POM": "POM",
    "sacrificial template": "牺牲模板",
    "mercury intrusion porosimetry": "压汞法",
    "MIP": "MIP",
    "krypton physisorption": "氪气物理吸附",
    "Brunauer-Emmett-Teller": "BET",
    "BET": "BET",
    "inductively coupled plasma": "电感耦合等离子体",
    "ICP": "ICP",
    "turnover frequency": "转换频率",
    "hydrogen production": "制氢",
    "green hydrogen": "绿氢",
    "water splitting": "水分解",
    "electrolyzer": "电解槽",
    "cathode": "阴极",
    "anode": "阳极",
    "diaphragm": "隔膜",
    "pore size distribution": "孔径分布",
    "specific surface area": "比表面积",
    "tensile strength": "拉伸强度",
    "Raman spectra": "拉曼光谱",
    "binding energy": "结合能",
    "valence band": "价带",
    "coordination number": "配位数",
    "Debye-Waller factor": "Debye-Waller因子",
    "oxygen vacancy": "氧空位",
    "rate-determining step": "速控步骤",
    "Volmer step": "Volmer步骤",
    "Heyrovsky step": "Heyrovsky步骤",
    "Tafel step": "Tafel步骤",
    "active sites": "活性位点",
    "degradation rate": "衰减速率",
    "cell voltage": "槽电压",
    "energy consumption": "能耗",
    "steam methane reforming": "蒸汽甲烷重整",
    "coal gasification": "煤气化",
}

def apply_glossary(text):
    """Replace English terms with Chinese equivalents using glossary."""
    # Sort by length (longest first) to avoid partial matches
    for en, cn in sorted(GLOSSARY.items(), key=lambda x: -len(x[0])):
        # Match whole words only for multi-word terms
        if " " in en:
            text = re.sub(re.escape(en), cn, text, flags=re.IGNORECASE)
    return text

# ============================================================
# Translation engine
# ============================================================
def translate_text(text, use_api=True):
    """Translate English text to Chinese."""
    if not text or not text.strip():
        return ""

    # First apply glossary protection
    text = apply_glossary(text)

    if use_api:
        try:
            # Split long text into chunks
            if len(text) > 4000:
                chunks = []
                sentences = text.replace('. ', '.|').replace('? ', '?|').replace('! ', '!|').split('|')
                chunk = ""
                for s in sentences:
                    if len(chunk) + len(s) < 4000:
                        chunk += s + ". "
                    else:
                        chunks.append(chunk)
                        chunk = s + ". "
                if chunk:
                    chunks.append(chunk)

                result = []
                for chunk in chunks:
                    translated = GoogleTranslator(source='en', target='zh-CN').translate(chunk)
                    result.append(translated)
                return " ".join(result)
            else:
                return GoogleTranslator(source='en', target='zh-CN').translate(text)
        except Exception as e:
            print(f"  Translation API error: {e}, using glossary-only mode")
            return text  # Already glossary-applied

    return text


def translate_section(title_en, paragraphs):
    """Translate a section with title and body paragraphs."""
    title_cn = translate_text(title_en)
    body_cn = []
    for p in paragraphs:
        if p.strip():
            body_cn.append(translate_text(p.strip()))
    return {"title_cn": title_cn, "body_cn": body_cn, "title_en": title_en, "body_en": paragraphs}


# ============================================================
# PDF Parsing
# ============================================================
def get_text_blocks(pdf_path):
    """Extract text blocks with positions from PDF."""
    doc = fitz.open(pdf_path)
    all_pages = []
    for i, page in enumerate(doc):
        blocks = page.get_text("dict")["blocks"]
        text_blocks = []
        for b in blocks:
            if b["type"] == 0:  # text block
                text_lines = []
                for line in b["lines"]:
                    line_text = "".join([span["text"] for span in line["spans"]])
                    text_lines.append(line_text)
                full_text = " ".join(text_lines)
                if full_text.strip():
                    text_blocks.append({
                        "bbox": list(b["bbox"]),
                        "text": full_text.strip(),
                    })
        # Sort by y then x
        text_blocks.sort(key=lambda b: (b["bbox"][1], b["bbox"][0]))
        all_pages.append(text_blocks)
    doc.close()
    return all_pages


def find_figure_captions(text_blocks):
    """Find figure caption blocks and their bounding boxes."""
    captions = []
    for block in text_blocks:
        t = block["text"]
        if re.search(r'Figure\s+\d+', t, re.IGNORECASE):
            captions.append(block)
    return captions


def render_page_region(pdf_path, page_num, bbox, dpi=RENDER_DPI):
    """Render a specific region of a PDF page to PNG."""
    doc = fitz.open(pdf_path)
    page = doc[page_num]

    # Convert to PDF points
    x0, y0, x1, y1 = bbox
    # Clip to page bounds
    rect = page.rect
    x0 = max(x0 - 5, rect.x0)
    y0 = max(y0 - 5, rect.y0)
    x1 = min(x1 + 5, rect.x1)
    y1 = min(y1 + 5, rect.y1)

    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, clip=fitz.Rect(x0, y0, x1, y1))
    doc.close()
    return pix.tobytes("png")


def render_full_page(pdf_path, page_num, dpi=RENDER_DPI):
    """Render a full PDF page to PNG."""
    doc = fitz.open(pdf_path)
    page = doc[page_num]
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat)
    doc.close()
    return pix.tobytes("png")


def save_image(data, filename):
    """Save image bytes to file."""
    path = os.path.join(IMAGES_DIR, filename)
    with open(path, "wb") as f:
        f.write(data)
    print(f"  Saved: {filename}")
    return path


# ============================================================
# Figure extraction strategy
# ============================================================

def extract_main_figures():
    """
    Extract figures from main paper.
    Strategy: For each page that contains a figure, render the figure region
    based on caption positions. If caption detection fails, render full page.
    """
    print("\n=== Extracting Main Paper Figures ===")
    doc = fitz.open(MAIN_PDF)
    figures = {}

    # Page-by-page figure extraction plan
    # Based on analysis of the JACS paper layout:
    #   Page 1: No figures (title/abstract/intro text only)
    #   Page 2: Figure 1 parts (top area)
    #   Page 3: Figure 1 continued (bottom) + Figure 2 parts
    #   Page 4: Figure 2 continued + Figure 3 parts
    #   Page 5: Figure 3 continued + Figure 4 parts
    #   Page 6: Figure 4 continued + Figure 5 parts
    #   Page 7: Figure 5 continued
    #   Page 8-10: Text + references

    # Extract full pages for each figure-containing page
    figure_pages = {
        "Figure1": [1, 2],      # Pages 2-3
        "Figure2": [2, 3],      # Pages 3-4
        "Figure3": [3, 4],      # Pages 4-5
        "Figure4": [4, 5],      # Pages 5-6
        "Figure5": [5, 6],      # Pages 6-7
    }

    for fig_name, pages in figure_pages.items():
        for pnum in pages:
            # Try to find caption on page to get figure bottom
            blocks = get_text_blocks(MAIN_PDF)[pnum]
            captions = find_figure_captions(blocks)

            # Check if this figure's caption is on this page
            fig_num = int(fig_name.replace("Figure", ""))
            matching_captions = [c for c in captions
                                if re.search(rf'Figure\s+{fig_num}\b', c["text"])]

            if matching_captions:
                # Use the first matching caption as figure bottom
                cap = matching_captions[0]
                caption_bbox = cap["bbox"]
                figure_bbox = [50, 50, 560, caption_bbox[1] - 5]
                data = render_page_region(MAIN_PDF, pnum, figure_bbox)
                save_image(data, f"{fig_name}_page{pnum+1}.png")
                figures[f"{fig_name}_page{pnum+1}"] = f"images/{fig_name}_page{pnum+1}.png"
            else:
                # No matching caption, render full page
                data = render_full_page(MAIN_PDF, pnum)
                save_image(data, f"{fig_name}_page{pnum+1}_full.png")
                figures[f"{fig_name}_page{pnum+1}_full"] = f"images/{fig_name}_page{pnum+1}_full.png"

    doc.close()
    return figures


def extract_si_figures():
    """
    Extract SI figures.
    Strategy: For each SI page with a figure, render the full page.
    """
    print("\n=== Extracting SI Figures ===")
    doc = fitz.open(SI_PDF)
    si_figures = {}

    # SI pages 1-6: Title + Methods (no figures)
    # SI pages 7+: Figures S1-S28 + Tables
    # Figure pages typically have "Figure S" in text

    for i in range(doc.page_count):
        page = doc[i]
        text = page.get_text()
        # Check if page contains a figure
        fig_match = re.search(r'Figure\s+S(\d+)', text)
        if fig_match:
            fig_num = fig_match.group(1)
            data = render_full_page(SI_PDF, i)
            fname = f"Figure_S{fig_num}_page{i+1}.png"
            save_image(data, fname)
            si_figures[f"Figure_S{fig_num}"] = f"images/{fname}"

        # Also check for tables
        table_match = re.search(r'Table\s+S(\d+)', text)
        if table_match:
            table_num = table_match.group(1)
            data = render_full_page(SI_PDF, i)
            fname = f"Table_S{table_num}_page{i+1}.png"
            save_image(data, fname)
            si_figures[f"Table_S{table_num}"] = f"images/{fname}"

    doc.close()
    return si_figures


# ============================================================
# Text extraction for translation
# ============================================================

def extract_main_text():
    """Extract structured text from main paper."""
    print("\n=== Extracting Main Paper Text ===")
    doc = fitz.open(MAIN_PDF)
    sections = []

    # Extract structured content
    full_text = ""
    for i in range(doc.page_count):
        page = doc[i]
        text = page.get_text("text")
        full_text += text + "\n"

    doc.close()

    # Parse sections
    # The paper structure: ABSTRACT, INTRODUCTION, RESULTS AND DISCUSSION, CONCLUSION
    lines = full_text.split('\n')
    current_section = "Header"
    current_text = []

    sections = {
        "title": "",
        "abstract": "",
        "introduction": "",
        "results_intro": "",
        "figure1_text": "",
        "figure2_text": "",
        "figure3_text": "",
        "figure4_text": "",
        "figure5_text": "",
        "conclusion": "",
    }

    # Extract title (first few non-empty lines)
    title_lines = []
    for line in lines[:20]:
        line = line.strip()
        if line and not line.startswith("Cite This") and not line.startswith("Read Online") \
                and not line.startswith("ACCESS") and not line.startswith("Article") \
                and not line.startswith("*") and not line.startswith("Supporting"):
            if len(line) > 30:
                title_lines.append(line)
            if len(title_lines) >= 2:
                break
    sections["title"] = " ".join(title_lines)

    # Extract abstract
    abstract_start = full_text.find("ABSTRACT:")
    intro_start = full_text.find("■INTRODUCTION")
    if abstract_start >= 0 and intro_start > abstract_start:
        sections["abstract"] = full_text[abstract_start+9:intro_start].strip()

    # Extract introduction
    results_start = full_text.find("■RESULTS AND DISCUSSION")
    if intro_start >= 0 and results_start > intro_start:
        sections["introduction"] = full_text[intro_start+12:results_start].strip()

    # Results sections - identify by Figure references
    results_text = ""
    conclusion_start = full_text.find("■CONCLUSION")
    if results_start >= 0:
        if conclusion_start > results_start:
            results_text = full_text[results_start+22:conclusion_start].strip()
        else:
            results_text = full_text[results_start+22:].strip()

    sections["results_intro"] = results_text[:1500]  # First portion

    # Find figure-specific text
    fig1_idx = results_text.find("Figure 1")
    fig2_idx = results_text.find("Figure 2")
    fig3_idx = results_text.find("Figure 3")
    fig4_idx = results_text.find("Figure 4")
    fig5_idx = results_text.find("Figure 5")

    for name, start, end in [
        ("figure1_text", fig1_idx, fig2_idx),
        ("figure2_text", fig2_idx, fig3_idx),
        ("figure3_text", fig3_idx, fig4_idx),
        ("figure4_text", fig4_idx, fig5_idx),
        ("figure5_text", fig5_idx, len(results_text)),
    ]:
        if start >= 0 and end > start:
            # Clean up caption references
            text = results_text[start:end]
            # Remove figure caption lines
            text = re.sub(r'Figure\s+\d+[\.\s].*?\n', '', text)
            sections[name] = text.strip()

    # Extract conclusion
    if conclusion_start >= 0:
        conclusion_text = full_text[conclusion_start+11:]
        # Stop at ASSOCIATED CONTENT
        assoc_idx = conclusion_text.find("■ASSOCIATED CONTENT")
        if assoc_idx > 0:
            conclusion_text = conclusion_text[:assoc_idx]
        sections["conclusion"] = conclusion_text.strip()

    return sections


def build_translation_draft():
    """Build comprehensive translation draft."""
    print("\n=== Building Translation Draft ===")

    # Key paragraphs to translate (manually curated)
    draft = {
        "slides": []
    }

    # ---- Slide metadata and curated text ----
    slide_data = [
        {
            "id": 1,
            "type": "title",
            "title_cn": "",
            "content_en": [],
            "content_cn": [],
            "figures": [],
        },
        # Will be populated below
    ]

    # Let's define the slide structure with English content to translate
    slides_content = [
        # Slide 1: Title
        {
            "id": 1, "type": "title",
            "title_en": "An Atomic-to-Macroscale Assembled Ni/MoO2 Electrode for High-Efficiency and Long-Life Hydrogen Production",
            "title_cn_curated": "原子-宏观多尺度组装Ni/MoO₂电极用于高效长寿命制氢",
            "subtitle_en": "Shang Jiang, Wei Hu, Shizheng Zhou, Linfeng Yu, Liang Luo, Yunlong Zhang, Qiao Zhao, Zhibin Yu, Wei Liu, Xiaoming Sun, Liang Yu*, Yanting Liu*, Dehui Deng*\nDalian Institute of Chemical Physics, CAS\nJ. Am. Chem. Soc. 2025",
            "subtitle_cn_curated": "蒋尚‡, 胡伟‡, 周世正, 于林峰, 罗亮, 张云龙, 赵翘, 于志彬, 刘伟, 孙晓明, 余亮*, 刘艳廷*, 邓德会*\n中国科学院大连化学物理研究所\nJ. Am. Chem. Soc. 2025",
            "note_cn": "",
        },
        # Slide 2: Research Background
        {
            "id": 2, "type": "content",
            "title_cn_curated": "研究背景与挑战",
            "bullets_en": [
                "Green hydrogen production via water electrolysis powered by renewable electricity is critical for establishing a green hydrogen economy.",
                "Alkaline water electrolysis (ALKWE) uses nonprecious metal catalysts and mature infrastructure but costs 3-4× more than steam methane reforming.",
                "Current ALKWE: high energy consumption (4.5-5.0 kWh/Nm³ H₂) and limited current density (200-400 mA/cm²).",
                "Key challenge: activity-stability-mass transport trade-off becomes severe at ampere-level current densities due to vigorous bubble evolution blocking active sites.",
            ],
            "bullets_cn_curated": [
                "利用可再生能源驱动的电化学水分解制氢是建立绿氢经济的关键过程。",
                "碱性水电解（ALKWE）可使用非贵金属催化剂和成熟工业基础设施，但成本比蒸汽甲烷重整高3-4倍。",
                "当前ALKWE面临高能耗（4.5-5.0 kWh/Nm³ H₂）和有限电流密度（200-400 mA/cm²）的挑战。",
                "核心难题：在高电流密度下，剧烈的气泡演化会阻塞活性位点、阻碍电解质进入并破坏催化层，使活性-稳定性-传质之间的矛盾尤为突出。",
            ],
        },
        # Slide 3: This Work Highlights
        {
            "id": 3, "type": "content",
            "title_cn_curated": "本文亮点：原子到宏观多尺度电极设计",
            "bullets_en": [
                "Atomic-to-macroscale assembly of integrated Ni/MoO₂ electrode with triscale (nano-micro-macro) porosity.",
                "Ultralow overpotential of 145 mV at 1 A/cm² (vs 300 mV for Pt/C) in 1 M KOH.",
                "Stable operation for >3500 h with degradation rate of 0.023 mV/h.",
                "Practical ALKWE: cell voltage 1.80 V at 1 A/cm², energy consumption 4.3 kWh/Nm³ H₂, meeting DOE 2026 target.",
                "Triple-enhancement effect: (i) electron transfer boosts intrinsic activity, (ii) hierarchical porosity enhances mass transfer, (iii) strong electronic interaction strengthens stability.",
            ],
            "bullets_cn_curated": [
                "原子到宏观多尺度组装：构建具有三尺度（纳-微-宏）孔结构的Ni/MoO₂一体化电极。",
                "超高活性：1 M KOH中，1 A/cm²下过电位仅145 mV（Pt/C为300 mV），降低52%。",
                "优异稳定性：在1 A/cm²下稳定运行超过3500小时，衰减速率仅0.023 mV/h。",
                "实际应用：工业条件下（30 wt% KOH, ≥85°C）槽电压1.80 V@1 A/cm²，能耗4.3 kWh/Nm³ H₂，达到美国DOE 2026年目标。",
                "三增强效应：(i) 界面电荷转移提升本征活性，(ii) 多级孔结构增强质量传输，(iii) 强电子相互作用增强结构稳定性。",
            ],
        },
        # Slide 4: Synthesis
        {
            "id": 4, "type": "figure_left_text_right",
            "title_cn_curated": "电极制备流程与结构表征（一）：合成示意图",
            "text_en": "The Ni/MoO₂ electrode is fabricated through powder metallurgy followed by impregnation-calcination: (1) Ni + NaCl powders sintered at 750°C in Ar, (2) NaCl template removed by water washing to create macropores (~400 μm), (3) impregnation in ammonium molybdate solution forming POM, (4) calcination at 600°C forming uniform Ni/MoO₂ heterostructure. The electrode retains triscale porosity with mass loading of 6.3 mg/cm².",
            "text_cn_curated": "Ni/MoO₂电极通过粉末冶金+浸渍-煅烧法制备：(1) Ni粉与NaCl模板混合压制成型，750°C Ar气氛烧结；(2) 水洗去除NaCl模板，形成~400 μm大孔；(3) 钼酸铵溶液中浸渍，表面Ni与溶液反应生成POM；(4) 600°C煅烧，形成均匀的Ni/MoO₂异质结构。电极保持三尺度孔结构，活性层载量6.3 mg/cm²。",
            "figures": ["Figure1_page2", "Figure_S1"],
        },
        # Slide 5: SEM morphology
        {
            "id": 5, "type": "figure_left_text_right",
            "title_cn_curated": "电极制备与结构表征（二）：微观形貌",
            "text_en": "Top-view SEM shows dual-scale porosity retained with enhanced surface roughness. Cross-sectional SEM reveals ~500 nm-thick Ni/MoO₂ layer uniformly coating the porous Ni skeleton. HAADF-STEM shows nanopores forming within the layer on micropore ligaments, with Ni nanoparticles anchored on nano-MoO₂.",
            "text_cn_curated": "顶视SEM显示双尺度孔结构保持完整，表面粗糙度增加。截面SEM显示约500 nm厚的Ni/MoO₂层均匀包裹多孔Ni骨架。HAADF-STEM显示纳米孔在微孔韧带上的生长层内形成，Ni纳米颗粒锚定在纳米MoO₂上。比表面积达2.91 m²/g，是商用Ni泡沫的100倍。",
            "figures": ["Figure1_page2", "Figure_S5"],
        },
        # Slide 6: Cross-section
        {
            "id": 6, "type": "figure_left_text_right",
            "title_cn_curated": "电极制备与结构表征（三）：截面结构与元素分布",
            "text_en": "FIB-prepared cross-section shows uniform Mo and Ni dispersion throughout the ~500 nm layer. HAADF-STEM and EDS mapping confirm abundant Ni nanoparticles anchored to nano-MoO₂ matrix. The hierarchical porous architecture is preserved after Ni/MoO₂ growth.",
            "text_cn_curated": "FIB制备的截面样品显示Mo和Ni元素在~500 nm层内均匀分布。HAADF-STEM和EDS元素面扫证实大量Ni纳米颗粒均匀锚定在纳米MoO₂基底上。Ni/MoO₂生长后保持了基体的多级孔结构。",
            "figures": ["Figure1_page3", "Figure_S6", "Figure_S8"],
        },
        # Slide 7: Atomic structure
        {
            "id": 7, "type": "figure_left_text_right",
            "title_cn_curated": "电极制备与结构表征（四）：原子尺度异质界面",
            "text_en": "HRTEM reveals d-spacings of 0.20 nm (Ni(111)) and 0.24 nm (MoO₂(200)), confirming the formation of Ni/MoO₂ heterostructure. EDS mapping shows clear spatial relationship between Ni nanoparticles and MoO₂ substrate. XRD and Raman confirm crystalline MoO₂ and metallic Ni phases.",
            "text_cn_curated": "HRTEM显示晶面间距0.20 nm（Ni(111)）和0.24 nm（MoO₂(200)），证实Ni/MoO₂异质结构的形成。EDS元素面扫显示Ni纳米颗粒与MoO₂基底的空间对应关系。XRD和拉曼光谱确认电极由晶态MoO₂和金属Ni相组成。",
            "figures": ["Figure1_page3", "Figure_S9", "Figure_S10", "Figure_S11"],
        },
        # Slide 8: Pore architecture
        {
            "id": 8, "type": "figure_left_text_right",
            "title_cn_curated": "电极制备与结构表征（五）：三尺度孔结构",
            "text_en": "MIP results show that Ni/MoO₂ growth leads to slight reduction in submillimeter- and micrometer-scale pores while introducing nanoporosity. The electrode possesses triscale (macro-micro-nano) porous architecture. Macro pores (~400 μm) from NaCl templates, micro pores from Ni powder sintering, and nano pores from the Ni/MoO₂ layer.",
            "text_cn_curated": "压汞法（MIP）结果显示Ni/MoO₂生长后亚毫米和微米孔略有减少，同时引入纳米孔。电极具有纳-微-宏三尺度孔结构：大孔（~400 μm）来自NaCl模板，微孔来自Ni粉烧结，纳米孔来自Ni/MoO₂层生长。这种多级孔结构协同增强电催化性能。",
            "figures": ["Figure1_page3_full", "Figure_S7"],
        },
        # Slide 9: LSV performance
        {
            "id": 9, "type": "figure_left_text_right",
            "title_cn_curated": "HER电化学性能（一）：极化曲线与活性对比",
            "text_en": "Ni/MoO₂ electrode exhibits overpotential of 19 mV at 10 mA/cm² (close to Pt/C at 14 mV). At 1 A/cm², the overpotential is only 145 mV, representing a 52% reduction compared to Pt/C (300 mV). Porous Ni requires 202 mV at 10 mA/cm². Tafel slope: 36 mV/dec for Ni/MoO₂ vs 87 mV/dec for porous Ni, indicating faster HER kinetics.",
            "text_cn_curated": "Ni/MoO₂电极在10 mA/cm²下过电位仅19 mV（接近Pt/C的14 mV），远优于多孔Ni（202 mV）。在1 A/cm²安培级电流密度下，过电位仅145 mV，比Pt/C（300 mV）降低52%。Tafel斜率：Ni/MoO₂为36 mV/dec，多孔Ni为87 mV/dec，表明更快的HER动力学。",
            "figures": ["Figure2_page3", "Figure_S12", "Figure_S13"],
        },
        # Slide 10: EIS
        {
            "id": 10, "type": "figure_left_text_right",
            "title_cn_curated": "HER电化学性能（二）：阻抗与电化学活性面积",
            "text_en": "EIS at -0.065 V vs RHE shows Ni/MoO₂ has much smaller charge transfer resistance (Rct = 0.426 Ω·cm²) compared to porous Ni (26.47 Ω·cm²). ECSA of Ni/MoO₂ is 211.25, which is 4.2× that of porous Ni (50.00), indicating more accessible active sites.",
            "text_cn_curated": "EIS测试显示Ni/MoO₂的电荷转移阻抗Rct仅0.426 Ω·cm²，远低于多孔Ni（26.47 Ω·cm²），表明更快的电荷转移动力学。电化学活性面积ECSA为211.25，是多孔Ni（50.00）的4.2倍，说明活性位点密度显著增加。",
            "figures": ["Figure2_page4_full"],
        },
        # Slide 11: Cdl and ECSA
        {
            "id": 11, "type": "content",
            "title_cn_curated": "HER电化学性能（三）：双电层电容与活性面积",
            "bullets_en": [],
            "bullets_cn_curated": [
                "通过CV在非法拉第电位窗口测定双电层电容（Cdl）。",
                "Ni/MoO₂的Cdl为8.45 mF/cm²，是多孔Ni（2.00 mF/cm²）的4.2倍。",
                "ECSA = Cdl/Cs = 211.25（Cs = 0.04 mF/cm²），多孔Ni ECSA = 50.00。",
                "更高的ECSA表明Ni/MoO₂电极提供了更多的可及活性位点。",
            ],
            "figures": ["Figure_S14"],
        },
        # Slide 12: H2 productivity
        {
            "id": 12, "type": "figure_left_text_right",
            "title_cn_curated": "HER电化学性能（四）：氢气产率与法拉第效率",
            "text_en": "The Ni/MoO₂ electrode shows H₂ productivity linearly increasing with applied current density, reaching near 100% Faradaic efficiency. The improved ECSA contributes to enhanced catalytic activity, enabling high current densities with nearly complete conversion efficiency.",
            "text_cn_curated": "Ni/MoO₂电极的氢气产率随外加电流密度线性增加，法拉第效率接近100%。提升的ECSA赋予多孔Ni更高的催化活性，使其能够实现高电流密度下近乎完全转化的产氢效率。",
            "figures": ["Figure2_page4"],
        },
        # Slide 13: Stability
        {
            "id": 13, "type": "figure_left_text_right",
            "title_cn_curated": "HER电化学性能（五）：长期稳定性测试",
            "text_en": "Ni/MoO₂ electrode maintains stable cathode potential at 1 A/cm² for over 3500 h with a low degradation rate of 0.023 mV/h. XRD, SEM show no obvious structural changes after the test. ICP confirms negligible Mo leaching (2.74 at% → 2.65 at%). Performance is competitive among nonprecious metal HER catalysts.",
            "text_cn_curated": "Ni/MoO₂电极在1 A/cm²下稳定运行超过3500小时，电位衰减速率仅0.023 mV/h。耐久性测试后XRD和SEM显示晶体结构和形貌无明显变化。ICP证实Mo含量几乎无损失（2.74 at% → 2.65 at%）。该性能在非贵金属碱性HER催化剂中具有极强的竞争力。",
            "figures": ["Figure2_page4_full", "Figure_S15", "Figure_S16"],
        },
        # Slide 14: ALKWE performance
        {
            "id": 14, "type": "figure_left_text_right",
            "title_cn_curated": "HER电化学性能（六）：实际ALKWE全电解测试",
            "text_en": "Ni/MoO₂ as cathode in practical ALKWE: cell voltage 1.69 V at 500 mA/cm² (4.0 kWh/Nm³ H₂). At 1 A/cm² and 95°C: cell voltage 1.80 V, energy consumption 4.3 kWh/Nm³ H₂ — meeting DOE 2026 target. Stable operation for >1000 h at 1 A/cm² with voltage degradation of 0.16 mV/h.",
            "text_cn_curated": "以Ni/MoO₂为阴极、Raney Ni为阳极构建ALKWE电解槽。工业条件下（30 wt% KOH, 85-95°C）：500 mA/cm²槽电压1.69 V，能耗4.0 kWh/Nm³ H₂；1 A/cm²下槽电压1.80 V，能耗4.3 kWh/Nm³ H₂，达到美国DOE 2026年目标。在1 A/cm²下稳定运行超1000 h，电压衰减0.16 mV/h。",
            "figures": ["Figure2_page4_full", "Figure_S17"],
        },
        # Slide 15: Performance comparison
        {
            "id": 15, "type": "content",
            "title_cn_curated": "性能对比：文献对标",
            "bullets_cn_curated": [
                "Table S6：与文献中非贵金属碱性HER催化剂对比，Ni/MoO₂在安培级电流密度（1 A/cm²）下的过电位（145 mV）和稳定性（3700 h, 0.023 mV/h）均处于领先水平。",
                "Table S8：与文献ALKWE电解槽性能对比，Ni/MoO₂||Raney Ni体系在1 A/cm²下1.80 V的槽电压是目前最优值之一，稳定性（1100 h）也远超同类体系（通常100 h）。",
                "综合性能达到甚至超过了DOE 2026年目标，展示了其工业应用潜力。",
            ],
            "figures": ["Table_S6", "Table_S8"],
        },
        # Slide 16: Bubble evolution
        {
            "id": 16, "type": "figure_left_text_right",
            "title_cn_curated": "质量传输能力（一）：气泡演化原位观测",
            "text_en": "In situ high-speed video imaging monitors H₂ bubble evolution at high current densities. At 500 mA/cm², Ni/MoO₂ produces bubbles with median diameter of 19 μm (IQR: 14-30 μm), markedly smaller than porous Ni (median: 29 μm, IQR: 20-43 μm). Smaller bubbles indicate enhanced detachment efficiency, reducing active site blockage.",
            "text_cn_curated": "原位高速摄像动态监测高电流密度下H₂气泡演化行为。500 mA/cm²下，Ni/MoO₂电极产生的气泡中位直径仅19 μm（IQR: 14-30 μm），显著小于多孔Ni（中位直径29 μm, IQR: 20-43 μm）。气泡尺寸的显著减小表明气泡脱附效率提高，有效减少了气体积累对活性位点的阻塞。",
            "figures": ["Figure3_page4"],
        },
        # Slide 17: Bubble size analysis
        {
            "id": 17, "type": "figure_left_text_right",
            "title_cn_curated": "质量传输能力（二）：气泡尺寸分布对比",
            "text_en": "The pronounced reduction in bubble size indicates enhanced bubble detachment efficiency, reducing the blockage of active sites and alleviating structural degradation caused by bubble-induced stress. At 1 A/cm², Ni/MoO₂ continues to generate smaller bubbles compared to porous Ni.",
            "text_cn_curated": "气泡尺寸的显著减小表明：(1) 提高了气泡脱附效率，减少气体积累对活性位点的阻塞；(2) 减轻了气泡诱导应力对催化层的结构破坏。在1 A/cm²更高电流密度下，Ni/MoO₂电极依然保持较小的气泡尺寸，展现出优异的气泡管理能力。",
            "figures": ["Figure3_page5_full", "Figure_S19"],
        },
        # Slide 18: Wettability
        {
            "id": 18, "type": "figure_left_text_right",
            "title_cn_curated": "质量传输能力（三）：超亲水性与快速浸润",
            "text_en": "Dynamic wettability tests show Ni/MoO₂ electrode facilitates complete electrolyte permeation within 50 ms, contrasting with porous Ni (contact angle 135°). DFT shows H₂O* adsorption is stronger on MoO₂ (ΔGads = -0.58 eV) than on Ni(111) (ΔGads = 0.02 eV), explaining enhanced hydrophilicity.",
            "text_cn_curated": "动态润湿性测试显示Ni/MoO₂电极在50 ms内即可实现完全电解质浸润，而多孔Ni呈现135°的接触角（疏水）。DFT计算揭示其本质原因：MoO₂(200)表面的H₂O*吸附能（ΔGads = -0.58 eV）远强于Ni(111)表面（ΔGads = 0.02 eV），MoO₂的引入赋予了电极增强的本征亲水性。",
            "figures": ["Figure3_page5"],
        },
        # Slide 19: Aerophobicity
        {
            "id": 19, "type": "figure_left_text_right",
            "title_cn_curated": "质量传输能力（四）：超疏气性与低粘附力",
            "text_en": "Ni/MoO₂ exhibits superior aerophobicity with H₂ bubble contact angle of 153° (vs 143° for porous Ni). Adhesion force measurements show 5-fold reduction: 0.58 μN for Ni/MoO₂ vs 2.9 μN for porous Ni. Nanoscale surface roughness contributes to both superhydrophilicity and superaerophobicity.",
            "text_cn_curated": "Ni/MoO₂电极展现出优异的疏气性：H₂气泡接触角153°（多孔Ni为143°）。气泡粘附力测量显示5倍降低：Ni/MoO₂仅0.58 μN，多孔Ni为2.9 μN。纳米级表面粗糙度同时贡献了超亲水性和超疏气性，实现了电解质快速渗透与气泡高效脱附的协同。",
            "figures": ["Figure3_page5_full"],
        },
        # Slide 20: Mass transfer summary
        {
            "id": 20, "type": "content",
            "title_cn_curated": "质量传输能力（五）：多级孔结构传质机理",
            "bullets_cn_curated": [
                "宏-微孔（~400 μm + ~5 μm）：提供快速电解质渗透通道，50 ms内完成浸润。",
                "纳米孔+亲水MoO₂涂层：增强毛细管驱动的液体渗透，保证活性位点处电解质供应。",
                "纳米表面粗糙度：赋予超疏气性（153°接触角），气泡粘附力降低5倍。",
                "协同效应：快速电解质渗透 + 高效气泡脱附 → 缓解高电流密度下的传质极化损失，保证长期稳定性。",
            ],
            "figures": [],
        },
        # Slide 21: XPS analysis
        {
            "id": 21, "type": "figure_left_text_right",
            "title_cn_curated": "电子结构分析（一）：XPS揭示Ni→MoO₂电子转移",
            "text_en": "XPS reveals: (i) Ni 2p₃/₂ Ni⁰ peak in Ni/MoO₂ shifts to higher binding energy vs porous Ni → electron-deficient Ni. (ii) Mo 3d Mo⁴⁺ and O 1s lattice oxygen peaks shift to lower binding energy vs pure MoO₂ → electron-enriched MoO₂. These observations confirm interfacial electron transfer from Ni to MoO₂.",
            "text_cn_curated": "XPS分析揭示界面电子转移：(i) Ni/MoO₂中Ni⁰的Ni 2p₃/₂峰相对于多孔Ni向高结合能位移→Ni呈电子缺失态；(ii) Mo 3d中Mo⁴⁺峰和O 1s中晶格氧峰相对于纯MoO₂向低结合能位移→MoO₂呈电子富集态。以上结果清晰证实了Ni→MoO₂的界面电子转移。",
            "figures": ["Figure4_page5", "Figure_S20"],
        },
        # Slide 22: Ni XANES/EXAFS
        {
            "id": 22, "type": "figure_left_text_right",
            "title_cn_curated": "电子结构分析（二）：Ni K-edge XANES与EXAFS",
            "text_en": "Ni K-edge XANES shows absorption edge shift to higher energy for Ni/MoO₂, confirming electron depletion at Ni sites. FT-EXAFS shows no detectable Ni-O scattering path (negligible nickel oxides). WT-EXAFS confirms dominant Ni-Ni metallic bonding contribution. Fitting: Ni-Ni CN = 4.8, R = 2.48 Å.",
            "text_cn_curated": "Ni K-edge XANES显示Ni/MoO₂吸收边向高能方向移动，确认Ni位点电子缺失。FT-EXAFS未检测到Ni-O散射路径，说明氧化镍信号可忽略。WT-EXAFS证实主要贡献来自Ni-Ni金属键。拟合结果：Ni-Ni配位数4.8，键长2.48 Å，以金属态Ni为主。",
            "figures": ["Figure4_page6", "Figure_S21", "Figure_S22"],
        },
        # Slide 23: Mo XANES
        {
            "id": 23, "type": "figure_left_text_right",
            "title_cn_curated": "电子结构分析（三）：Mo K-edge XANES",
            "text_en": "Mo K-edge XANES confirms Mo predominantly retains +4 oxidation state in Ni/MoO₂. The absorption edge position is consistent with MoO₂ reference, indicating the Mo valence state is preserved after heterostructure formation despite interfacial electron transfer.",
            "text_cn_curated": "Mo K-edge XANES确认Ni/MoO₂中Mo主要保持+4氧化态。吸收边位置与MoO₂参比一致，表明尽管存在界面电子转移，Mo的价态在异质结构形成后基本保持不变。MoO₂作为电子受体增强了Ni位点的催化活性。",
            "figures": ["Figure4_page6_full"],
        },
        # Slide 24: Mo EXAFS
        {
            "id": 24, "type": "figure_left_text_right",
            "title_cn_curated": "电子结构分析（四）：Mo K-edge EXAFS - 局域结构畸变",
            "text_en": "Mo K-edge FT-EXAFS reveals distorted local structure in Ni/MoO₂ compared to pure MoO₂: Mo-O bond slightly elongated (2.02 vs 2.00 Å) and Mo-O coordination number significantly reduced (4.6 vs 5.9). Distortion attributed to oxygen vacancies (from reducing synthesis atmosphere) and interfacial electron transfer from Ni.",
            "text_cn_curated": "Mo K-edge FT-EXAFS揭示Ni/MoO₂中Mo局域结构相对纯MoO₂发生畸变：Mo-O键略微伸长（2.02 vs 2.00 Å），Mo-O配位数显著降低（4.6 vs 5.9）。结构畸变归因于：(1) 合成过程中还原气氛产生的氧空位；(2) Ni→MoO₂界面电子转移诱导的额外电子/结构扰动。",
            "figures": ["Figure4_page6_full", "Figure_S23", "Figure_S24"],
        },
        # Slide 25: Electronic structure summary
        {
            "id": 25, "type": "content",
            "title_cn_curated": "电子结构分析（五）：界面电子转移机制总结",
            "bullets_cn_curated": [
                "XPS + XANES + EXAFS多维度证实：Ni/MoO₂异质界面存在Ni→MoO₂的定向电子转移。",
                "Ni位点：电子缺失→优化H*吸附能→提升本征HER活性。",
                "MoO₂位点：电子富集+局域结构畸变（氧空位+配位不饱和）→增强亲水性和导电性。",
                "强Ni-MoO₂电子相互作用→增强结构稳定性，耐受高电流密度下长期运行。",
            ],
            "figures": [],
        },
        # Slide 26: In situ XANES
        {
            "id": 26, "type": "figure_left_text_right",
            "title_cn_curated": "反应机理（一）：原位XANES - 工作状态下的电子结构",
            "text_en": "In situ XANES under HER conditions reveals the electronic structure of Ni and Mo remains largely unchanged during catalysis. Ni K-edge and Mo K-edge spectra at different applied potentials (HER times) overlap with open-circuit spectra, indicating high structural stability of Ni/MoO₂ heterointerface under working conditions.",
            "text_cn_curated": "原位XANES实验在HER工作条件下监测电子结构演化。不同外加电位（不同反应时间）下采集的Ni K-edge和Mo K-edge谱图与开路电位谱图几乎完全重合，表明在HER工作条件下Ni/MoO₂异质界面的电子结构保持高度稳定，这是长期催化稳定性的结构基础。",
            "figures": ["Figure5_page6", "Figure_S25"],
        },
        # Slide 27: DFT models
        {
            "id": 27, "type": "figure_left_text_right",
            "title_cn_curated": "反应机理（二）：DFT计算 - 界面模型与差分电荷",
            "text_en": "DFT models: Ni(111) surface, MoO₂(200) surface, and Ni/MoO₂ heterointerface (Ni(111) nanoribbon on MoO₂(200)). Differential charge distribution shows pronounced charge redistribution at the interface: blue (charge depletion) on Ni side, yellow (charge accumulation) on MoO₂ side, consistent with XPS/XANES results.",
            "text_cn_curated": "DFT计算构建了Ni(111)、MoO₂(200)和Ni/MoO₂界面（Ni(111)纳米带置于MoO₂(200)表面）模型。差分电荷分布显示界面处显著的电荷重新分布：Ni侧呈蓝色（电荷耗散），MoO₂侧呈黄色（电荷积累），与XPS/XANES实验结论一致，证实Ni→MoO₂界面电子转移。",
            "figures": ["Figure5_page7_full", "Figure_S26", "Figure_S27"],
        },
        # Slide 28: H* adsorption
        {
            "id": 28, "type": "figure_left_text_right",
            "title_cn_curated": "反应机理（三）：DFT计算 - H*吸附自由能优化",
            "text_en": "PDOS shows electron-deficient Ni at the interface shifts d-band center, weakening H* adsorption. Ni/MoO₂ interface Ni site has ΔGH* = -0.18 eV, much closer to 0 eV than pure Ni(111) (-0.43 eV). This moderate H* binding facilitates both H* adsorption (Volmer step) and H₂ desorption, thereby boosting overall HER kinetics.",
            "text_cn_curated": "PDOS分析显示界面处电子缺失的Ni其d带中心发生移动，适度减弱了H*吸附强度。界面Ni位点的ΔGH* = -0.18 eV，比纯Ni(111)的-0.43 eV更接近热中性（0 eV）。适中的H*结合能同时有利于H*吸附（Volmer步骤）和H₂脱附，从而显著提升整体HER动力学。",
            "figures": ["Figure5_page7"],
        },
        # Slide 29: Reaction pathway
        {
            "id": 29, "type": "figure_left_text_right",
            "title_cn_curated": "反应机理（四）：反应路径与三增强效应协同机制",
            "text_en": "The triple-enhancement mechanism: (i) Ni→MoO₂ electron transfer weakens H* adsorption on Ni sites (ΔGH* from -0.43 to -0.18 eV), boosting intrinsic activity; (ii) Triscale porosity with hydrophilic MoO₂ accelerates bubble detachment and electrolyte permeation, enhancing mass transfer; (iii) Strong Ni-MoO₂ electronic interaction ensures structural stability under high current densities.",
            "text_cn_curated": "三增强效应协同机制：(i) 本征活性增强：Ni→MoO₂界面电子转移优化Ni位点H*吸附自由能（ΔGH*从-0.43优化至-0.18 eV），促进H₂脱附；(ii) 传质增强：三尺度孔结构+亲水MoO₂涂层协同加速气泡脱附和电解质渗透；(iii) 稳定性增强：强Ni-MoO₂电子相互作用和与电极骨架的牢固结合确保高电流密度下结构稳定。",
            "figures": ["Figure5_page7_full"],
        },
        # Slide 30: Conclusion
        {
            "id": 30, "type": "content",
            "title_cn_curated": "总结与展望",
            "bullets_cn_curated": [
                "创新点：提出了原子到宏观多尺度组装策略，成功构建了具有三尺度孔结构的Ni/MoO₂一体化电极。",
                "核心性能：1 A/cm²下过电位145 mV（较Pt/C降低52%），3500+小时稳定运行，工业条件下槽电压1.80 V。",
                "三增强效应：界面电荷转移→本征活性↑；多级孔+亲水涂层→传质↑；强电子作用→稳定性↑。",
                "制备优势：粉末冶金+浸渍-煅烧法，原料廉价易得，工艺与现有工业制造兼容，NaCl模板可回收。",
                "展望：该多尺度电极设计策略为克服高电流密度ALKWE中的活性-稳定性权衡提供了新思路。",
            ],
            "figures": [],
        },
        # Slide 31: Key data summary
        {
            "id": 31, "type": "content",
            "title_cn_curated": "关键性能数据汇总",
            "bullets_cn_curated": [
                "过电位 @10 mA/cm²：19 mV（Pt/C: 14 mV）",
                "过电位 @1 A/cm²：145 mV（Pt/C: 300 mV，降低52%）",
                "Tafel斜率：36 mV/dec（多孔Ni: 87 mV/dec）",
                "电荷转移阻抗Rct：0.426 Ω·cm²（多孔Ni: 26.47 Ω·cm²）",
                "ECSA：211.25（多孔Ni: 50.00，4.2倍）",
                "法拉第效率：~100%",
                "稳定性：3500+ h @1 A/cm²，衰减0.023 mV/h",
                "ALKWE槽电压：1.80 V @1 A/cm²（DOE 2026目标）",
                "ALKWE能耗：4.3 kWh/Nm³ H₂",
                "ALKWE稳定性：1000+ h，电压衰减0.16 mV/h",
            ],
            "figures": [],
        },
        # Slide 32: Acknowledgment
        {
            "id": 32, "type": "title",
            "title_cn_curated": "谢谢！欢迎提问",
            "subtitle_en": "",
            "subtitle_cn_curated": "Reference: Jiang, Hu, et al. J. Am. Chem. Soc. 2025, https://doi.org/10.1021/jacs.5c21735",
            "note_cn": "",
        },
    ]

    # Translate all English content
    print("Translating slide content...")
    for i, slide in enumerate(slides_content):
        print(f"  Slide {slide['id']}...")
        # Translate title if no curated translation
        if "title_cn_curated" not in slide or not slide["title_cn_curated"]:
            if "title_en" in slide and slide["title_en"]:
                slide["title_cn"] = translate_text(slide["title_en"])
            else:
                slide["title_cn"] = slide.get("title_cn_curated", "")

        # Translate bullets
        if "bullets_en" in slide and slide["bullets_en"]:
            if "bullets_cn_curated" in slide and slide["bullets_cn_curated"]:
                slide["bullets_cn"] = slide["bullets_cn_curated"]
            else:
                slide["bullets_cn"] = [translate_text(b) for b in slide["bullets_en"]]

        # Translate text
        if "text_en" in slide and slide["text_en"]:
            if "text_cn_curated" in slide and slide["text_cn_curated"]:
                slide["text_cn"] = slide["text_cn_curated"]
            else:
                slide["text_cn"] = translate_text(slide["text_en"])

        # Translate subtitle
        if "subtitle_en" in slide and slide["subtitle_en"]:
            if "subtitle_cn_curated" in slide and slide["subtitle_cn_curated"]:
                slide["subtitle_cn"] = slide["subtitle_cn_curated"]
            else:
                slide["subtitle_cn"] = translate_text(slide["subtitle_en"])

    draft["slides"] = slides_content
    return draft


# ============================================================
# PPT Builder
# ============================================================

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

FONT_TITLE = 'Microsoft YaHei'
FONT_BODY = 'Microsoft YaHei'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_bar(slide, title_text, slide_num=None):
    """Add a dark blue title bar at the top of the slide."""
    # Title bar background
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE

    # Slide number
    if slide_num:
        num_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
        nf = num_box.text_frame
        np = nf.paragraphs[0]
        np.text = str(slide_num)
        np.font.size = Pt(9)
        np.font.color.rgb = LIGHT_GRAY
        np.font.name = FONT_BODY
        np.alignment = PP_ALIGN.RIGHT


def add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_BODY
    p.alignment = alignment
    p.line_spacing = line_spacing
    return tf


def add_bullet_points(slide, left, top, width, height, bullets, font_size=Pt(14)):
    """Add bullet points to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = font_size
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_BODY
        p.line_spacing = Pt(22)
        p.space_after = Pt(6)
    return tf


def resolve_image_path(fig_name):
    """Find the actual image file given a figure reference name."""
    # Try multiple possible paths
    candidates = [
        f"D:/firstcc/images/{fig_name}",
        f"D:/firstcc/images/{fig_name}.png",
        f"D:/firstcc/{fig_name}",
        f"D:/firstcc/{fig_name}.png",
    ]
    # Also search in images dir for files containing the fig_name
    # Use exact prefix match: fig_name + "_" or fig_name + "." to avoid partial matches
    # e.g., "Figure_S1" matches "Figure_S1_page7.png" but not "Figure_S10_page19.png"
    if os.path.isdir("D:/firstcc/images"):
        for f in sorted(os.listdir("D:/firstcc/images")):
            if f == fig_name + ".png" or f.startswith(fig_name + "_") or f.startswith(fig_name + "."):
                candidates.insert(0, f"D:/firstcc/images/{f}")
    for c in candidates:
        if os.path.exists(c):
            return c
    return None


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add image to slide, handling missing files gracefully."""
    full_path = resolve_image_path(os.path.basename(img_path))
    if full_path is None:
        full_path = resolve_image_path(img_path)
    if full_path and os.path.exists(full_path):
        try:
            if width:
                slide.shapes.add_picture(full_path, left, top, width=width)
            elif height:
                slide.shapes.add_picture(full_path, left, top, height=height)
            else:
                slide.shapes.add_picture(full_path, left, top)
            return True
        except Exception as e:
            print(f"  Warning: Could not add image {full_path}: {e}")
    return False


def build_ppt(slides_data, output_path):
    """Build the PowerPoint presentation."""
    print("\n=== Building PPT ===")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # blank

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # Set background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        slide_id = slide_info["id"]
        slide_type = slide_info.get("type", "content")

        # Get title (prefer curated)
        title = slide_info.get("title_cn_curated") or slide_info.get("title_cn", "")
        add_title_bar(slide, title, slide_id)

        if slide_type == "title":
            # Title slide - centered text
            subtitle = slide_info.get("subtitle_cn_curated") or slide_info.get("subtitle_cn", "")
            note = slide_info.get("note_cn", "")
            if subtitle:
                add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(3.5),
                            subtitle, Pt(14), color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                            line_spacing=1.8)
            if note:
                add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(1.0),
                            note, Pt(10), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        elif slide_type == "content":
            # Content slide - bullets only
            bullets = slide_info.get("bullets_cn_curated") or slide_info.get("bullets_cn", [])
            figs = slide_info.get("figures", [])

            if bullets:
                bullet_left = Inches(0.8)
                bullet_width = Inches(7.0) if figs else Inches(11.5)
                add_bullet_points(slide, bullet_left, Inches(1.3), bullet_width, Inches(5.8),
                                bullets, Pt(14))

            # Add figures on the right if any
            if figs:
                fig_top = Inches(1.3)
                fig_left = Inches(8.2)
                fig_w = Inches(4.5)
                fig_height = Inches(5.5 / max(len(figs), 1))
                for j, fig in enumerate(figs):
                    fpath = resolve_image_path(fig)
                    if fpath:
                        try:
                            slide.shapes.add_picture(fpath, fig_left,
                                                    fig_top + Inches(j * 5.5 / len(figs)),
                                                    width=fig_w)
                        except Exception:
                            pass

        elif slide_type == "figure_left_text_right":
            # Figure on left, text on right
            text_cn = slide_info.get("text_cn_curated") or slide_info.get("text_cn", "")
            figs = slide_info.get("figures", [])

            if figs:
                fig_left = Inches(0.5)
                fig_top = Inches(1.3)
                fig_width = Inches(5.8)
                fig_height = Inches(5.5 / max(len(figs), 1))
                for j, fig in enumerate(figs):
                    fpath = resolve_image_path(fig)
                    if fpath:
                        try:
                            slide.shapes.add_picture(fpath, fig_left,
                                                    fig_top + Inches(j * 5.5 / len(figs)),
                                                    width=fig_width)
                        except Exception:
                            pass

            if text_cn:
                add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.8),
                            text_cn, Pt(13), color=DARK_GRAY, line_spacing=1.8)

        # Add source note at bottom
        add_text_box(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
                    "Jiang, Hu, et al. J. Am. Chem. Soc. 2025", Pt(7), color=LIGHT_GRAY)

    prs.save(output_path)
    print(f"\nPPT saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


# ============================================================
# Main
# ============================================================
def stage1():
    """Stage 1: Extract images and generate translation draft."""
    print("=" * 60)
    print("STAGE 1: Extracting images & generating translations")
    print("=" * 60)

    # Ensure images directory exists
    os.makedirs(IMAGES_DIR, exist_ok=True)

    # Extract figures from main paper
    main_figs = extract_main_figures()
    print(f"\nMain paper figures extracted: {len(main_figs)}")

    # Extract SI figures
    si_figs = extract_si_figures()
    print(f"SI figures extracted: {len(si_figs)}")

    # Build translation draft
    draft = build_translation_draft()

    # Save translation draft
    with open(TRANSLATION_FILE, "w", encoding="utf-8") as f:
        json.dump(draft, f, ensure_ascii=False, indent=2)
    print(f"\nTranslation draft saved to: {TRANSLATION_FILE}")
    print(f"Total slides defined: {len(draft['slides'])}")

    print("\n" + "=" * 60)
    print("STAGE 1 COMPLETE")
    print(f"Please review and edit: {TRANSLATION_FILE}")
    print(f"Review images in: {IMAGES_DIR}")
    print("Then run: python ppt_generator.py stage2")
    print("=" * 60)


def stage2():
    """Stage 2: Build PPT from reviewed translations."""
    print("=" * 60)
    print("STAGE 2: Building PPT from reviewed translations")
    print("=" * 60)

    # Load reviewed translations
    if not os.path.exists(TRANSLATION_FILE):
        print(f"ERROR: {TRANSLATION_FILE} not found. Run stage1 first.")
        sys.exit(1)

    with open(TRANSLATION_FILE, "r", encoding="utf-8") as f:
        draft = json.load(f)

    slides_data = draft.get("slides", [])
    if not slides_data:
        print("ERROR: No slides found in translation file.")
        sys.exit(1)

    # Build PPT
    build_ppt(slides_data, OUTPUT_PPT)

    print("\n" + "=" * 60)
    print("STAGE 2 COMPLETE")
    print(f"PPT saved to: {OUTPUT_PPT}")
    print("=" * 60)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage:")
        print("  python ppt_generator.py stage1   # Extract images + translations")
        print("  python ppt_generator.py stage2   # Build PPT from reviewed translations")
        sys.exit(1)

    stage = sys.argv[1].lower()
    if stage == "stage1":
        stage1()
    elif stage == "stage2":
        stage2()
    else:
        print(f"Unknown stage: {stage}")
        print("Use 'stage1' or 'stage2'")
        sys.exit(1)
